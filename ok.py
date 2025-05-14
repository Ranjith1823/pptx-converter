from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
import sqlite3
import copy
from collections import defaultdict
import random
from datetime import datetime

def safe_append_rows(table, needed_count):
    existing = len(table.rows)
    if existing == 0:
        print("⚠️ Table has no rows to clone!")
        return
    base_row = table.rows[1] if len(table.rows) > 1 else table.rows[0]
    while len(table.rows) < needed_count:
        table._tbl.append(copy.deepcopy(base_row._tr))

def duplicate_slide(prs, slide):
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide.shapes:
        new_shape = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    return new_slide

# Connect and fetch data between dates

# Ask for date range
from_date = input("Enter the start date (YYYY-MM-DD): ")
to_date = input("Enter the end date (YYYY-MM-DD): ")

# Validate date format
try:
    from_date = datetime.strptime(from_date, "%Y-%m-%d").date()
    to_date = datetime.strptime(to_date, "%Y-%m-%d").date()
except ValueError:
    raise ValueError("Invalid date format. Please use YYYY-MM-DD.")

# Ensure from_date is not after to_date
if from_date > to_date:
    raise ValueError("Start date cannot be after end date.")

# Connect to the database and fetch data
conn = sqlite3.connect('db.sqlite3', detect_types=sqlite3.PARSE_DECLTYPES | sqlite3.PARSE_COLNAMES)
cursor = conn.cursor()
cursor.execute("""
    SELECT 
        What_I_accomplished_at_work, 
        One_thing_I_Learnt, 
        Plan_for_the_next_2_weeks, 
        Support_at_school_or_shifu_or_at_work, 
        Insight, 
        Blog_Url ,
        Created_date
    FROM Notes_notes
    WHERE DATE(Created_date) BETWEEN ? AND ?
""", (from_date, to_date))
data = cursor.fetchall()
conn.close()

# Organize data
work_accomplished = defaultdict(list)
things_learned = []
plans = defaultdict(list)
support_needed = []
insights = []
blog_urls = []

for record in data:
    work_entry = record[0].strip()
    learn_entry = record[1].strip()
    plan_entry = record[2].strip()
    support_entry = record[3].strip()
    insight_entry = record[4].strip()
    blog_entry = record[5].strip()

    if ":" in work_entry:
        project_id, status = work_entry.split(":", 1)
        work_accomplished[project_id.strip()].append(status.strip())

    if learn_entry:
        things_learned.append(learn_entry)

    if ":" in plan_entry:
        project_id, plan_status = plan_entry.split(":", 1)
        plans[project_id.strip()].append(plan_status.strip())

    if support_entry:
        support_needed.append(support_entry)

    if insight_entry:
        insights.append(insight_entry)

    if blog_entry.startswith("http"):
        blog_urls.append(blog_entry)

# Load template
prs = Presentation("C3SLD_Weekly_19_10_2024(model).pptx")

# Slide 1: Date
for shape in prs.slides[0].shapes:
    if shape.has_text_frame and "Weekly status" in shape.text:
        for p in shape.text_frame.paragraphs:
            if "@" in p.text:
                parts = p.text.split("@")
                p.text = f"{parts[0]}@  Auroville \t\t May-2025"

# Slide 3: Learning + Status + First 3 Project IDs
proj_items = list(work_accomplished.items())  # Initialize proj_items with project data
base_slide3 = prs.slides[2]

slide3 = base_slide3
for shape in slide3.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        col_count = len(table.columns)
        if table.cell(0, col_count - 1).text.lower().startswith("learn"):
            sampled_learnings = random.sample(things_learned, min(8, len(things_learned)))  # Randomly select up to 8 learnings
            safe_append_rows(table, len(sampled_learnings) + 1)
            for i, item in enumerate(sampled_learnings):
                table.cell(i + 1, col_count - 1).text = item

        if "project" in table.cell(0, 0).text.lower():
            first_chunk = proj_items[:3]  # First 3 project IDs
            safe_append_rows(table, len(first_chunk) + 1)
            for j, (proj, statuses) in enumerate(first_chunk):
                table.cell(j + 1, 0).text = proj
                table.cell(j + 1, 1).text = "\n- ".join([""] + statuses)  # Start each status with "-"

# Slide 4: Remaining Project Work + Teach — paginate
base_slide4 = prs.slides[3]
proj_items = proj_items[3:]  # Exclude the first 3 project IDs already used in Slide 3
teach_items = support_needed
chunk_size = 3  # Limit to 3 project IDs per slide
teach_chunk_size = 10  # Limit to 10 teach items per slide

for i in range(0, len(proj_items), chunk_size):
    chunk = proj_items[i:i + chunk_size]
    slide = base_slide4 if i == 0 else duplicate_slide(prs, base_slide4)
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            if "project" in table.cell(0, 0).text.lower():
                safe_append_rows(table, len(chunk) + 1)
                for j, (proj, statuses) in enumerate(chunk):
                    table.cell(j + 1, 0).text = proj
                    table.cell(j + 1, 1).text = "\n- ".join([""] + statuses)  # Start each status with "-"

            elif "teach" in table.cell(0, 0).text.lower():
                teach_start = i * teach_chunk_size
                teach_end = teach_start + teach_chunk_size
                teach_chunk = teach_items[teach_start:teach_end]
                safe_append_rows(table, len(teach_chunk) + 1)
                for j, item in enumerate(teach_chunk):
                    table.cell(j + 1, 0).text = item

# Slide 5: Project Plan — paginate
base_slide5 = prs.slides[4]
project_status = [{"project": k, "status": v} for k, v in plans.items()]
for i in range(0, len(project_status), chunk_size):
    chunk = project_status[i:i + chunk_size]
    slide = base_slide5 if i == 0 else duplicate_slide(prs, base_slide5)
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            safe_append_rows(table, len(chunk) + 1)
            for j, item in enumerate(chunk):
                table.cell(j + 1, 0).text = item["project"]
                table.cell(j + 1, 1).text = "\n- ".join([""] + item["status"])  # Start each status with "-"

# Slide 6: Insights + Blogs
base_slide6 = prs.slides[5]
slide = base_slide6
target_shape = None
for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip() == "":
        target_shape = shape
        break

if not target_shape:
    raise Exception("Target shape not found on Slide 6")

tf = target_shape.text_frame
tf.clear()

# Add Insights
p = tf.paragraphs[0]
p.text = "Insights:"
p.level = 0
p.font.bold = True
p.font.size = Pt(18)
for insight in random.sample(insights, min(5, len(insights))):
    para = tf.add_paragraph()
    para.text = insight
    para.level = 1
    para.font.size = Pt(14)

# Add Blogs
tf.add_paragraph()
p = tf.add_paragraph()
p.text = "Blogs:"
p.level = 0
p.font.bold = True
p.font.size = Pt(18)
p._element.get_or_add_pPr().insert(
    0, parse_xml(r'<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
)
for blog in blog_urls:
    para = tf.add_paragraph()
    para.text = blog
    para.level = 1
    para.font.size = Pt(14)

# Save presentation
prs.save("C3SLD_Final_AutoPaginated.pptx")
print("✅ Presentation generated: C3SLD_Final_AutoPaginated.pptx")
